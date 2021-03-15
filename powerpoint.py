#
# Written by Ramtin Mesgari
# For use by Blacfox
#

from pptx import Presentation

##############
# # FIELDS # #
##############

file = "Blacfox Template.pptx"
prs = Presentation(file)

slides = prs.slides  # get slides
slide_index = []  # slide index
slide_idx = []  # slide id

# PLACEHOLDERS #
place_holders = []

heading = None
introduction = None
before_heading = None
before_text = None
middle_heading = None
middle_text = None
after_heading = None
after_text = None
quote = None
percentage1 = None
percentage_text1 = None
percentage2 = None
percentage_text2 = None
customer_name_heading = None
customer_name_text = None
industry_heading = None
industry_text = None
products_and_services_heading = None
products_and_services_text = None
employees_heading = None
employees_text = None
revenue_heading = None
revenue_text = None
featured_solutions_heading = None
featured_solutions_text = None
video_heading = None
video_text = None


#################
# # FUNCTIONS # #
#################

def get_slide(index):
    return slides.get(slide_idx[slide_index[index]])


def run():
    global place_holders

    for slide in slides:
        slide_index.append(slides.index(slide))
        slide_idx.append(slides[slides.index(slide)].slide_id)

    slide = prs.slides.add_slide(prs.slide_layouts[0])

    for shape in slide.shapes:
        place_holders.append(shape)

    print(len(place_holders))

    for i in range(len(place_holders)):
        place_holders[i].text = "Placeholder: " + str(i + 1)
        print("\t|\t", place_holders[i].text, "\t|")

    create_placeholders()


def create_placeholders():
    global heading
    global introduction
    global before_heading
    global before_text
    global middle_heading
    global middle_text
    global after_heading
    global after_text
    global quote
    global percentage1
    global percentage_text1
    global percentage2
    global percentage_text2
    global customer_name_heading
    global customer_name_text
    global industry_heading
    global industry_text
    global products_and_services_heading
    global products_and_services_text
    global employees_heading
    global employees_text
    global revenue_heading
    global revenue_text
    global featured_solutions_heading
    global featured_solutions_text
    global video_heading
    global video_text

    heading = place_holders[27]
    introduction = place_holders[20]
    before_heading = place_holders[26]
    before_text = place_holders[25]
    middle_heading = place_holders[24]
    middle_text = place_holders[23]
    after_heading = place_holders[22]
    after_text = place_holders[21]
    quote = place_holders[19]
    percentage1 = place_holders[18]
    percentage_text1 = place_holders[17]
    percentage2 = place_holders[16]
    percentage_text2 = place_holders[15]
    customer_name_heading = place_holders[14]
    customer_name_text = place_holders[13]
    industry_heading = place_holders[1]
    industry_text = place_holders[0]
    products_and_services_heading = place_holders[3]
    products_and_services_text = place_holders[2]
    employees_heading = place_holders[5]
    employees_text = place_holders[4]
    revenue_heading = place_holders[7]
    revenue_text = place_holders[6]
    featured_solutions_heading = place_holders[9]
    featured_solutions_text = place_holders[8]
    video_heading = place_holders[11]
    video_text = place_holders[10]


def save(new_file):
    prs.save(new_file)
